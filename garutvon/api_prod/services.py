from __future__ import annotations

import base64
import gzip
import io
import re
from collections import Counter
from pathlib import Path
from typing import Any

from fastapi import HTTPException
from PIL import Image, UnidentifiedImageError


def _extract_text_from_bytes(content: bytes, filename: str, content_type: str | None = None) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix in {".txt", ".md", ".csv", ".log"} or content_type == "text/plain":
        return content.decode("utf-8", errors="ignore")

    if suffix == ".pdf" or content_type == "application/pdf":
        try:
            import pdfplumber

            text = []
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    text.append(page.extract_text() or "")
            return "\n".join(text).strip()
        except Exception:
            pass

    if suffix in {".docx", ".doc"} or content_type in {"application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"}:
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs if p.text)
        except Exception:
            pass

    if suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff"} or (content_type and content_type.startswith("image/")):
        return _ocr_image(content)

    if suffix == ".pptx" or content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            lines: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        lines.append(shape.text)
            return "\n".join(lines)
        except Exception:
            pass

    return content.decode("utf-8", errors="ignore")


def _summarize_text(text: str, max_sentences: int = 3) -> str:
    text = text.strip()
    if not text:
        return ""
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", text) if s.strip()]
    if len(sentences) <= max_sentences:
        return " ".join(sentences)

    words = [w.lower() for w in re.findall(r"\w+", text) if len(w) > 3]
    frequency = Counter(words)
    sentence_scores = {
        sentence: sum(frequency[word.lower()] for word in re.findall(r"\w+", sentence) if len(word) > 3)
        for sentence in sentences
    }
    selected = sorted(sentences, key=lambda s: sentence_scores[s], reverse=True)[:max_sentences]
    selected.sort(key=lambda s: sentences.index(s))
    return " ".join(selected)


def _ocr_image(content: bytes) -> str:
    try:
        import easyocr

        image = Image.open(io.BytesIO(content)).convert("RGB")
        reader = easyocr.Reader(["en"], gpu=False)
        result = reader.readtext(image)
        return "\n".join([item[1] for item in result])
    except Exception:
        try:
            import pytesseract

            image = Image.open(io.BytesIO(content)).convert("RGB")
            return pytesseract.image_to_string(image)
        except Exception:
            raise HTTPException(
                status_code=501,
                detail="OCR requires easyocr or pytesseract. Install one of those libraries to enable this service.",
            )


def summarize(content: str, target_language: str | None = None) -> dict[str, Any]:
    if not content:
        raise HTTPException(status_code=400, detail="No text provided for summarization.")

    summary = _summarize_text(content)
    return {"summary": summary, "sentence_count": len([s for s in re.split(r'(?<=[.!?])\s+', content) if s.strip()])}


def extract_text(content: bytes, filename: str, content_type: str | None = None) -> dict[str, Any]:
    text = _extract_text_from_bytes(content, filename, content_type)
    return {"filename": filename, "bytes": len(content), "text": text}


def metadata(content: bytes, filename: str, content_type: str | None = None) -> dict[str, Any]:
    result: dict[str, Any] = {
        "filename": filename,
        "bytes": len(content),
        "content_type": content_type or "application/octet-stream",
    }

    suffix = Path(filename).suffix.lower()
    if suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff"} or (content_type and content_type.startswith("image/")):
        try:
            image = Image.open(io.BytesIO(content))
            result.update({
                "format": image.format,
                "width": image.width,
                "height": image.height,
                "mode": image.mode,
            })
        except UnidentifiedImageError:
            pass

    if suffix == ".pdf" or content_type == "application/pdf":
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(content))
            result["pages"] = len(reader.pages)
            metadata = reader.metadata
            if metadata:
                result["metadata"] = {k: v for k, v in metadata.items() if v}
        except Exception:
            pass

    if suffix in {".docx", ".doc"} or content_type in {"application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"}:
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            result["paragraph_count"] = len(doc.paragraphs)
        except Exception:
            pass

    try:
        import magic

        result["detected_mime"] = magic.from_buffer(content, mime=True)
    except Exception:
        pass

    return result


def _base64_payload(content: bytes) -> str:
    return base64.b64encode(content).decode("utf-8")


def convert(content: bytes, filename: str, target_format: str, content_type: str | None = None) -> dict[str, Any]:
    target = target_format.lower().strip().lstrip(".")
    source_text = _extract_text_from_bytes(content, filename, content_type)
    output_bytes = b""
    output_filename = f"{Path(filename).stem}.{target}"

    if target == "txt":
        output_bytes = source_text.encode("utf-8")
    elif target == "pdf":
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas

            buffer = io.BytesIO()
            page = canvas.Canvas(buffer, pagesize=letter)
            text_object = page.beginText(40, 760)
            for line in source_text.splitlines():
                text_object.textLine(line[:100])
                if text_object.getY() < 40:
                    page.drawText(text_object)
                    page.showPage()
                    text_object = page.beginText(40, 760)
            page.drawText(text_object)
            page.save()
            output_bytes = buffer.getvalue()
        except Exception:
            raise HTTPException(status_code=500, detail="PDF conversion requires reportlab.")
    elif target == "docx":
        try:
            from docx import Document

            document = Document()
            for line in source_text.splitlines():
                document.add_paragraph(line)
            buffer = io.BytesIO()
            document.save(buffer)
            output_bytes = buffer.getvalue()
        except Exception:
            raise HTTPException(status_code=500, detail="DOCX conversion requires python-docx.")
    elif target in {"png", "jpg", "jpeg", "bmp", "gif"}:
        try:
            image = Image.open(io.BytesIO(content))
            if target == "jpg":
                target = "jpeg"
            buffer = io.BytesIO()
            image.save(buffer, format=target.upper())
            output_bytes = buffer.getvalue()
        except Exception:
            raise HTTPException(status_code=500, detail="Image conversion requires a valid image file.")
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported target format: {target_format}")

    return {
        "filename": output_filename,
        "target_format": target,
        "original_bytes": len(content),
        "converted_bytes": len(output_bytes),
        "output_base64": _base64_payload(output_bytes),
    }


def compress(content: bytes, filename: str, content_type: str | None = None) -> dict[str, Any]:
    suffix = Path(filename).suffix.lower()
    if suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff"} or (content_type and content_type.startswith("image/")):
        try:
            image = Image.open(io.BytesIO(content)).convert("RGB")
            buffer = io.BytesIO()
            image.save(buffer, format="JPEG", quality=65, optimize=True)
            compressed_bytes = buffer.getvalue()
        except Exception:
            compressed_bytes = gzip.compress(content)
    else:
        compressed_bytes = gzip.compress(content)

    return {
        "filename": filename,
        "original_bytes": len(content),
        "compressed_bytes": len(compressed_bytes),
        "compression_ratio": round(len(compressed_bytes) / len(content), 4) if content else 0.0,
        "compressed_base64": _base64_payload(compressed_bytes),
    }


def background_remove(content: bytes, filename: str, content_type: str | None = None) -> dict[str, Any]:
    try:
        from rembg import remove

        result = remove(content)
        return {
            "filename": f"{Path(filename).stem}_transparent.png",
            "original_bytes": len(content),
            "result_bytes": len(result),
            "output_base64": _base64_payload(result),
            "message": "Background removed successfully.",
        }
    except Exception:
        raise HTTPException(
            status_code=501,
            detail="Background removal requires rembg to be installed and available.",
        )
